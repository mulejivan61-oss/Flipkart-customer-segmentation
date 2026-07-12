
"""
=========================================================
  Customer Segmentation Intelligence Dashboard  (v2)
  Built with: Streamlit + scikit-learn + Plotly
  Supports: ANY dataset (numeric + categorical, any column
            names), multiple clustering algorithms, auto-K
            suggestion, deeper insights, optional AI summary
            (Anthropic API), and PDF / PPTX report export.
=========================================================
Run with:
    pip install streamlit pandas numpy scikit-learn plotly xlsxwriter openpyxl \
                matplotlib reportlab python-pptx anthropic
    streamlit run app.py
 
Note: matplotlib / reportlab / python-pptx / anthropic are only needed for the
"Reports" tab (PDF/PPTX export, AI summary). The rest of the dashboard works
fine without them.
"""
 
import streamlit as st
import pandas as pd
import numpy as np
from io import BytesIO
from sklearn.preprocessing import StandardScaler, OneHotEncoder
from sklearn.compose import ColumnTransformer
from sklearn.cluster import KMeans, AgglomerativeClustering, DBSCAN
from sklearn.mixture import GaussianMixture
from sklearn.metrics import silhouette_score
import plotly.express as px
import plotly.graph_objects as go
 
# Optional dependencies — dashboard still works if these are missing
try:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    MATPLOTLIB_OK = True
except ImportError:
    MATPLOTLIB_OK = False
 
try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    from reportlab.lib import colors as rl_colors
    from reportlab.platypus import (SimpleDocTemplate, Table, TableStyle, Paragraph,
                                     Spacer, Image as RLImage)
    from reportlab.lib.styles import getSampleStyleSheet
    REPORTLAB_OK = True
except ImportError:
    REPORTLAB_OK = False
 
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_OK = True
except ImportError:
    PPTX_OK = False
 
try:
    import anthropic
    ANTHROPIC_OK = True
except ImportError:
    ANTHROPIC_OK = False
 
# =========================================================
# PAGE CONFIG
# =========================================================
st.set_page_config(
    page_title="Customer Segmentation Dashboard",
    page_icon="🛍️",
    layout="wide",
    initial_sidebar_state="expanded"
)
 
# =========================================================
# PREMIUM DARK THEME CSS
# =========================================================
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Poppins:wght@300;400;500;600;700;800&display=swap');
 
    html, body, [class*="css"]  {
        font-family: 'Poppins', sans-serif;
    }
 
    .stApp {
        background: radial-gradient(circle at 10% 0%, #1a1c2e 0%, #0e0f1a 45%, #0a0b12 100%);
        color: #eaeaf0;
    }
 
    section[data-testid="stSidebar"] {
        background: linear-gradient(180deg, #14162a 0%, #0d0e1a 100%);
        border-right: 1px solid rgba(255,255,255,0.06);
    }
 
    h1, h2, h3, h4 {
        color: #f5f5ff !important;
        font-weight: 700 !important;
    }
 
    .hero {
        padding: 26px 32px;
        border-radius: 20px;
        background: linear-gradient(120deg, rgba(124,58,237,0.25), rgba(56,189,248,0.15));
        border: 1px solid rgba(255,255,255,0.08);
        margin-bottom: 24px;
        box-shadow: 0 10px 40px rgba(80, 40, 200, 0.15);
    }
    .hero h1 {
        font-size: 2.1rem;
        margin: 0;
        background: linear-gradient(90deg, #a78bfa, #38bdf8, #34d399);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
    }
    .hero p {
        color: #b7b9d0;
        margin-top: 6px;
        font-size: 0.95rem;
    }
 
    .kpi-card {
        border-radius: 18px;
        padding: 18px 20px;
        background: linear-gradient(145deg, rgba(255,255,255,0.05), rgba(255,255,255,0.02));
        border: 1px solid rgba(255,255,255,0.08);
        box-shadow: 0 8px 24px rgba(0,0,0,0.25);
        transition: transform 0.2s ease;
    }
    .kpi-card:hover { transform: translateY(-4px); }
    .kpi-label {
        font-size: 0.78rem;
        color: #9ca0c0;
        text-transform: uppercase;
        letter-spacing: 0.06em;
        margin-bottom: 6px;
    }
    .kpi-value {
        font-size: 1.7rem;
        font-weight: 800;
        color: #ffffff;
    }
    .kpi-sub { font-size: 0.75rem; color: #7de3a3; margin-top: 4px;}
 
    .badge-high {
        background: rgba(34,197,94,0.15); color:#4ade80; padding:4px 12px;
        border-radius:999px; font-weight:600; font-size:0.85rem; border:1px solid rgba(74,222,128,0.4);
    }
    .badge-medium {
        background: rgba(250,204,21,0.15); color:#facc15; padding:4px 12px;
        border-radius:999px; font-weight:600; font-size:0.85rem; border:1px solid rgba(250,204,21,0.4);
    }
    .badge-low {
        background: rgba(248,113,113,0.15); color:#f87171; padding:4px 12px;
        border-radius:999px; font-weight:600; font-size:0.85rem; border:1px solid rgba(248,113,113,0.4);
    }
    .badge-noise {
        background: rgba(148,163,184,0.15); color:#94a3b8; padding:4px 12px;
        border-radius:999px; font-weight:600; font-size:0.85rem; border:1px solid rgba(148,163,184,0.4);
    }
 
    .insight-box {
        background: rgba(124,58,237,0.08);
        border-left: 4px solid #a78bfa;
        padding: 14px 18px;
        border-radius: 10px;
        margin-bottom: 12px;
        font-size: 0.92rem;
        color: #dcdcf0;
    }
 
    div.stButton > button {
        background: linear-gradient(90deg, #7c3aed, #38bdf8);
        color: white; border: none; border-radius: 10px; font-weight: 600;
        padding: 0.5rem 1.2rem;
    }
    div.stButton > button:hover { opacity: 0.9; }
 
    .stTabs [data-baseweb="tab-list"] { gap: 6px; }
    .stTabs [data-baseweb="tab"] {
        background: rgba(255,255,255,0.04);
        border-radius: 10px 10px 0 0;
        padding: 8px 18px;
        color: #b7b9d0;
    }
    .stTabs [aria-selected="true"] {
        background: rgba(124,58,237,0.25) !important;
        color: #fff !important;
    }
</style>
""", unsafe_allow_html=True)
 
# =========================================================
# HELPERS
# =========================================================
@st.cache_data
def load_data(file):
    if file.name.lower().endswith(".csv"):
        return pd.read_csv(file)
    return pd.read_excel(file)
 
 
def to_excel_bytes(dataframe):
    output = BytesIO()
    with pd.ExcelWriter(output, engine="xlsxwriter") as writer:
        dataframe.to_excel(writer, index=False, sheet_name="Clustered_Customers")
    return output.getvalue()
 
 
BADGE_COLORS = {
    "High Value Customer": "#4ade80",
    "Medium Value Customer": "#facc15",
    "Low Value Customer": "#f87171",
    "Noise / Outlier": "#94a3b8",
}
 
 
def badge_html(customer_type):
    if customer_type == "High Value Customer":
        return "<span class='badge-high'>🟢 High Value</span>"
    elif customer_type == "Medium Value Customer":
        return "<span class='badge-medium'>🟡 Medium Value</span>"
    elif customer_type == "Noise / Outlier":
        return "<span class='badge-noise'>⚪ Noise / Outlier</span>"
    else:
        return "<span class='badge-low'>🔴 Low Value</span>"
 
 
def assign_customer_type(sorted_cluster_index):
    """Generic Low/Medium/High labelling that works for any K, based on
    relative rank of cluster mean value (ascending). Excludes noise (-1)."""
    n = len(sorted_cluster_index)
    mapping = {}
    for i, idx in enumerate(sorted_cluster_index):
        pct = i / (n - 1) if n > 1 else 1.0
        if pct < 0.34:
            mapping[idx] = "Low Value Customer"
        elif pct < 0.67:
            mapping[idx] = "Medium Value Customer"
        else:
            mapping[idx] = "High Value Customer"
    return mapping
 
 
def kpi_card(label, value, sub=""):
    st.markdown(f"""
        <div class="kpi-card">
            <div class="kpi-label">{label}</div>
            <div class="kpi-value">{value}</div>
            <div class="kpi-sub">{sub}</div>
        </div>
    """, unsafe_allow_html=True)
 
 
def safe_mode(series):
    m = series.mode()
    return m.iloc[0] if not m.empty else "NA"
 
 
def feature_kpi_value(dataframe, col, is_numeric):
    if is_numeric:
        return f"{dataframe[col].mean():,.1f}", "average"
    else:
        top = safe_mode(dataframe[col])
        return f"{top}", "most common"
 
 
def fit_clustering(algorithm, X, k=None, eps=0.5, min_samples=5):
    """Generic fit wrapper — returns (model_or_None, labels)."""
    if algorithm == "KMeans":
        model = KMeans(n_clusters=k, random_state=42, n_init=10)
        labels = model.fit_predict(X)
    elif algorithm == "Agglomerative (Hierarchical)":
        model = AgglomerativeClustering(n_clusters=k)
        labels = model.fit_predict(X)
    elif algorithm == "Gaussian Mixture":
        model = GaussianMixture(n_components=k, random_state=42)
        labels = model.fit_predict(X)
    else:  # DBSCAN
        model = DBSCAN(eps=eps, min_samples=min_samples)
        labels = model.fit_predict(X)
    return model, labels
 
 
def scan_k(algorithm, X, k_range):
    """Try each k in k_range, return list of (k, silhouette_or_None, wcss_or_None)."""
    results = []
    for k in k_range:
        try:
            model, lbls = fit_clustering(algorithm, X, k=k)
        except Exception:
            continue
        wcss_val = getattr(model, "inertia_", None)
        sil = None
        if len(set(lbls)) >= 2:
            try:
                sil = silhouette_score(X, lbls)
            except Exception:
                sil = None
        results.append((k, sil, wcss_val))
    return results
 
 
# =========================================================
# HERO HEADER
# =========================================================
st.markdown("""
<div class="hero">
    <h1>🛍️ Customer Segmentation Intelligence Dashboard</h1>
    <p>Upload ANY customer data (numeric or categorical columns, any names) — pick an algorithm,
    auto-suggest the best K, and get deep, exportable insights on who your best customers are.</p>
</div>
""", unsafe_allow_html=True)
 
# =========================================================
# SIDEBAR — UPLOAD
# =========================================================
st.sidebar.markdown("## 📂 Upload Data")
uploaded_file = st.sidebar.file_uploader("Upload CSV / Excel", type=["csv", "xlsx", "xls"])
 
if uploaded_file is None:
    st.info("👈 Upload a CSV or Excel file from the sidebar to get started. "
            "Any column names work — numeric columns (like spend, orders) "
            "and categorical columns (like city, gender, plan type) are both supported.")
    st.stop()
 
df = load_data(uploaded_file)
 
# =========================================================
# SIDEBAR — CLEANING & FEATURE SELECTION
# =========================================================
st.sidebar.markdown("---")
st.sidebar.markdown("## 🧹 Data Cleaning")
 
rows_before = df.shape[0]
df = df.drop_duplicates()
dup_removed = rows_before - df.shape[0]
 
numeric_cols = df.select_dtypes(include=np.number).columns.tolist()
datetime_cols = df.select_dtypes(include=["datetime64[ns]", "datetimetz"]).columns.tolist()
categorical_cols = [c for c in df.columns if c not in numeric_cols and c not in datetime_cols]
 
usable_cols = numeric_cols + categorical_cols
if len(usable_cols) < 2:
    st.error("Your file needs at least 2 usable columns (numeric or categorical) to run clustering.")
    st.stop()
 
default_features = [c for c in ["Annual_spending", "Orders_count"] if c in df.columns]
if len(default_features) < 2:
    default_features = numeric_cols[:2] if len(numeric_cols) >= 2 else usable_cols[:2]
 
features = st.sidebar.multiselect(
    "Select features for clustering (numeric or categorical)",
    usable_cols,
    default=default_features
)
 
if len(features) < 2:
    st.warning("⚠️ Please select at least 2 features to continue.")
    st.stop()
 
numeric_features = [f for f in features if f in numeric_cols]
categorical_features = [f for f in features if f in categorical_cols]
 
high_card_warns = [f for f in categorical_features if df[f].nunique(dropna=True) > 50]
if high_card_warns:
    st.sidebar.warning(
        f"⚠️ High-cardinality column(s): {', '.join(high_card_warns)} (50+ unique values). "
        f"Clustering still runs, but consider excluding these for cleaner segments."
    )
 
missing_filled = 0
for f in numeric_features:
    missing_filled += df[f].isnull().sum()
    df[f] = df[f].fillna(df[f].mean())
for f in categorical_features:
    missing_filled += df[f].isnull().sum()
    fill_val = safe_mode(df[f]) if not df[f].dropna().empty else "Unknown"
    df[f] = df[f].fillna(fill_val).astype(str)
 
st.sidebar.success(f"✔ Removed {dup_removed} duplicate rows")
st.sidebar.success(f"✔ Filled {int(missing_filled)} missing values (mean for numeric, mode for categorical)")
 
id_col_options = ["(row index)"] + [c for c in df.columns]
id_col = st.sidebar.selectbox("🔍 Column to use as Customer ID / Name (for search)", id_col_options)
 
# =========================================================
# ENCODING + SCALING (any mix of numeric/categorical)
# =========================================================
transformers = []
if numeric_features:
    transformers.append(("num", StandardScaler(), numeric_features))
if categorical_features:
    transformers.append(("cat", OneHotEncoder(handle_unknown="ignore"), categorical_features))
 
preprocessor = ColumnTransformer(transformers)
X_scaled = preprocessor.fit_transform(df[features])
if hasattr(X_scaled, "toarray"):
    X_scaled = X_scaled.toarray()
 
# =========================================================
# SIDEBAR — MODEL SETTINGS (algorithm choice + auto-K)
# =========================================================
st.sidebar.markdown("---")
st.sidebar.markdown("## 🔢 Model Settings")
 
algorithm = st.sidebar.selectbox(
    "Clustering Algorithm",
    ["KMeans", "Agglomerative (Hierarchical)", "Gaussian Mixture", "DBSCAN (auto cluster count)"]
)
 
eps, min_samples = None, None
 
if algorithm != "DBSCAN (auto cluster count)":
    elbow_max_k = st.sidebar.slider("Max K to scan", 3, 12, 6)
 
    if "k_value" not in st.session_state:
        st.session_state["k_value"] = 3
 
    if st.sidebar.button("🎯 Auto-suggest Best K"):
        scan_results = scan_k(algorithm, X_scaled, range(2, elbow_max_k + 1))
        valid = [(k, s) for k, s, _ in scan_results if s is not None]
        if valid:
            best_k, best_score = max(valid, key=lambda x: x[1])
            st.session_state["k_value"] = best_k
            st.sidebar.success(f"Suggested K = {best_k} (silhouette = {best_score:.3f})")
        else:
            st.sidebar.warning("Could not compute a suggestion for this data/algorithm.")
 
    st.session_state["k_value"] = min(max(st.session_state["k_value"], 2), 10)
    k_value = st.sidebar.slider("Number of Clusters (K)", 2, 10, key="k_value")
else:
    elbow_max_k = 6
    k_value = None
    eps = st.sidebar.slider("DBSCAN: eps (neighborhood radius)", 0.1, 5.0, 0.8, 0.1)
    min_samples = st.sidebar.slider("DBSCAN: min_samples", 2, 20, 5)
 
# =========================================================
# FINAL MODEL FIT (generic across algorithms)
# =========================================================
final_model, clusters = fit_clustering(algorithm, X_scaled, k=k_value, eps=eps, min_samples=min_samples)
df["Cluster"] = clusters
 
noise_present = -1 in set(clusters)
valid_labels = [c for c in set(clusters) if c != -1]
 
try:
    if noise_present:
        mask = df["Cluster"] != -1
        score = silhouette_score(X_scaled[mask.values], df.loc[mask, "Cluster"]) if mask.sum() > 1 and len(set(df.loc[mask, "Cluster"])) > 1 else 0.0
    elif len(valid_labels) >= 2:
        score = silhouette_score(X_scaled, clusters)
    else:
        score = 0.0
except Exception:
    score = 0.0
 
# Cluster summary — numeric via mean, categorical via most common value (noise excluded)
valid_df = df[df["Cluster"] != -1]
summary_parts = []
if numeric_features:
    summary_parts.append(valid_df.groupby("Cluster")[numeric_features].mean().round(2))
if categorical_features:
    summary_parts.append(valid_df.groupby("Cluster")[categorical_features].agg(safe_mode))
summary = pd.concat(summary_parts, axis=1) if summary_parts else pd.DataFrame(index=valid_labels)
 
if numeric_features:
    rank_feature = numeric_features[0]
    sort_basis = valid_df.groupby("Cluster")[rank_feature].mean()
else:
    rank_feature = None
    sort_basis = valid_df["Cluster"].value_counts()
 
sorted_clusters = sort_basis.sort_values()
customer_type_map = assign_customer_type(sorted_clusters.index)
df["Customer_Type"] = df["Cluster"].map(customer_type_map).fillna("Noise / Outlier")
 
primary_feature = numeric_features[0] if numeric_features else features[0]
primary_is_numeric = primary_feature in numeric_features
secondary_feature = (numeric_features[1] if len(numeric_features) > 1
                      else (features[1] if len(features) > 1 and features[1] != primary_feature else features[0]))
secondary_is_numeric = secondary_feature in numeric_features
 
 
def predict_new_point(new_scaled_row):
    """Predict cluster for a new point, generic across algorithms."""
    if algorithm in ("KMeans", "Gaussian Mixture"):
        return int(final_model.predict(new_scaled_row)[0])
    # Agglomerative / DBSCAN have no native predict -> nearest centroid fallback
    centroids = {c: X_scaled[df["Cluster"].values == c].mean(axis=0) for c in valid_labels}
    if not centroids:
        return valid_labels[0] if valid_labels else -1
    dists = {c: np.linalg.norm(new_scaled_row[0] - cent) for c, cent in centroids.items()}
    return min(dists, key=dists.get)
 
 
# =========================================================
# SIDEBAR — FILTERS
# =========================================================
st.sidebar.markdown("---")
st.sidebar.markdown("## 🎛️ Filters")
type_options = ["High Value Customer", "Medium Value Customer", "Low Value Customer"]
if noise_present:
    type_options.append("Noise / Outlier")
type_filter = st.sidebar.multiselect("Customer Type", options=type_options, default=type_options)
 
range_filters = {}
category_filters = {}
for f in features:
    if f in numeric_features:
        fmin, fmax = float(df[f].min()), float(df[f].max())
        if fmin == fmax:
            fmax = fmin + 1
        range_filters[f] = st.sidebar.slider(f"{f} range", fmin, fmax, (fmin, fmax))
    else:
        options = sorted(df[f].dropna().unique().tolist(), key=lambda x: str(x))
        category_filters[f] = st.sidebar.multiselect(f"{f} values", options, default=options)
 
filtered_df = df[df["Customer_Type"].isin(type_filter)].copy()
for f, (lo, hi) in range_filters.items():
    filtered_df = filtered_df[(filtered_df[f] >= lo) & (filtered_df[f] <= hi)]
for f, selected_vals in category_filters.items():
    filtered_df = filtered_df[filtered_df[f].isin(selected_vals)]
 
# =========================================================
# KPI CARDS
# =========================================================
st.markdown("### 📊 Key Metrics")
c1, c2, c3, c4, c5 = st.columns(5)
with c1:
    kpi_card("Total Customers", f"{len(df):,}")
with c2:
    val, sub = feature_kpi_value(df, primary_feature, primary_is_numeric)
    kpi_card(f"{'Avg' if primary_is_numeric else 'Top'} {primary_feature}", val, sub)
with c3:
    val, sub = feature_kpi_value(df, secondary_feature, secondary_is_numeric)
    kpi_card(f"{'Avg' if secondary_is_numeric else 'Top'} {secondary_feature}", val, sub)
with c4:
    kpi_card("Silhouette Score", f"{score:.3f}", "closer to 1 = better")
with c5:
    high_pct = (df["Customer_Type"] == "High Value Customer").mean() * 100
    kpi_card("High Value %", f"{high_pct:.1f}%", "of total customers")
 
st.markdown("<br>", unsafe_allow_html=True)
 
# =========================================================
# RULE-BASED INSIGHTS (plain text, reused by UI + PDF/PPT)
# =========================================================
def generate_insight_lines():
    total = len(df)
    high_n = (df["Customer_Type"] == "High Value Customer").sum()
    med_n = (df["Customer_Type"] == "Medium Value Customer").sum()
    low_n = (df["Customer_Type"] == "Low Value Customer").sum()
    noise_n = (df["Customer_Type"] == "Noise / Outlier").sum()
 
    lines = []
    if primary_is_numeric:
        high_val = df.loc[df["Customer_Type"] == "High Value Customer", primary_feature].sum()
        total_val = df[primary_feature].sum()
        high_share = (high_val / total_val * 100) if total_val else 0
        lines.append(f"{high_n} customers ({high_n/total*100:.1f}%) are High Value, "
                      f"contributing {high_share:.1f}% of total {primary_feature}.")
    else:
        top_val = safe_mode(df.loc[df["Customer_Type"] == "High Value Customer", primary_feature]) if high_n else "N/A"
        lines.append(f"{high_n} customers ({high_n/total*100:.1f}%) are High Value; "
                      f"their most common {primary_feature} is {top_val}.")
 
    lines.append(f"{med_n} customers ({med_n/total*100:.1f}%) are Medium Value — a strong target for upsell/cross-sell.")
    lines.append(f"{low_n} customers ({low_n/total*100:.1f}%) are Low Value — consider re-engagement or win-back offers.")
    if noise_present:
        lines.append(f"{noise_n} customers ({noise_n/total*100:.1f}%) were flagged as Noise/Outliers by DBSCAN — "
                      f"unusual behaviour worth a manual look.")
 
    quality = ("a strong, well-separated segmentation." if score > 0.5 else
               "a reasonable segmentation; try adjusting K, features, or algorithm for tighter clusters." if score > 0.25 else
               "clusters overlap significantly; consider different features, K, or algorithm.")
    lines.append(f"Clustering quality (Silhouette Score) is {score:.3f} — {quality}")
 
    # distinguishing features per cluster (z-score based)
    if numeric_features:
        overall_mean = df[numeric_features].mean()
        overall_std = df[numeric_features].std().replace(0, 1)
        for cluster_id in sorted(valid_labels):
            cluster_rows = df[df["Cluster"] == cluster_id]
            ctype = customer_type_map.get(cluster_id, "Unknown")
            z = ((cluster_rows[numeric_features].mean() - overall_mean) / overall_std).abs().sort_values(ascending=False)
            top_feats = z.head(2).index.tolist()
            desc = ", ".join([f"{f} is {'above' if cluster_rows[f].mean() > overall_mean[f] else 'below'} average"
                               for f in top_feats])
            lines.append(f"Segment '{ctype}' (cluster {cluster_id}, n={len(cluster_rows)}): {desc}.")
    return lines
 
 
# =========================================================
# TABS
# =========================================================
tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs([
    "📈 Elbow & Clusters", "🎨 Explore & Search", "👤 Predict a Customer",
    "📋 Data & Downloads", "📊 Business Insights", "📄 Reports"
])
 
# ---------------------------------------------------------
# TAB 1 — ELBOW / SILHOUETTE + CLUSTER PLOT
# ---------------------------------------------------------
with tab1:
    col1, col2 = st.columns(2)
 
    with col1:
        if algorithm == "KMeans":
            st.markdown("#### 📈 Elbow Method (WCSS)")
            scan_results = scan_k("KMeans", X_scaled, range(1, elbow_max_k + 1))
            k_list = [r[0] for r in scan_results]
            wcss = [r[2] for r in scan_results]
            elbow_fig = go.Figure()
            elbow_fig.add_trace(go.Scatter(
                x=k_list, y=wcss, mode="lines+markers",
                line=dict(color="#a78bfa", width=3),
                marker=dict(size=9, color="#38bdf8", line=dict(width=2, color="#7c3aed"))
            ))
            elbow_fig.add_vline(x=k_value, line_dash="dash", line_color="#4ade80",
                                 annotation_text=f"Chosen K={k_value}", annotation_font_color="#4ade80")
            elbow_fig.update_layout(
                template="plotly_dark", paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
                xaxis_title="Number of Clusters (K)", yaxis_title="WCSS (Inertia)",
                margin=dict(l=10, r=10, t=10, b=10), height=380
            )
            st.plotly_chart(elbow_fig, use_container_width=True)
        elif algorithm != "DBSCAN (auto cluster count)":
            st.markdown("#### 📈 Silhouette Score vs K")
            scan_results = scan_k(algorithm, X_scaled, range(2, elbow_max_k + 1))
            k_list = [r[0] for r in scan_results if r[1] is not None]
            sils = [r[1] for r in scan_results if r[1] is not None]
            sil_fig = go.Figure()
            sil_fig.add_trace(go.Scatter(
                x=k_list, y=sils, mode="lines+markers",
                line=dict(color="#a78bfa", width=3),
                marker=dict(size=9, color="#38bdf8", line=dict(width=2, color="#7c3aed"))
            ))
            if k_value in k_list:
                sil_fig.add_vline(x=k_value, line_dash="dash", line_color="#4ade80",
                                   annotation_text=f"Chosen K={k_value}", annotation_font_color="#4ade80")
            sil_fig.update_layout(
                template="plotly_dark", paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
                xaxis_title="Number of Clusters (K)", yaxis_title="Silhouette Score",
                margin=dict(l=10, r=10, t=10, b=10), height=380
            )
            st.plotly_chart(sil_fig, use_container_width=True)
        else:
            st.markdown("#### 📈 DBSCAN Summary")
            n_clusters_found = len(valid_labels)
            noise_count = int((df["Cluster"] == -1).sum())
            st.metric("Clusters found", n_clusters_found)
            st.metric("Noise points", f"{noise_count} ({noise_count/len(df)*100:.1f}%)")
            st.caption("DBSCAN decides the number of clusters automatically based on eps/min_samples — "
                       "tune those in the sidebar to get more or fewer clusters.")
 
    with col2:
        st.markdown("#### ⭐ Silhouette Score (final model)")
        gauge = go.Figure(go.Indicator(
            mode="gauge+number",
            value=round(score, 3),
            number={'font': {'color': '#ffffff'}},
            gauge={
                'axis': {'range': [-1, 1], 'tickcolor': '#ffffff'},
                'bar': {'color': "#a78bfa"},
                'bgcolor': "rgba(0,0,0,0)",
                'steps': [
                    {'range': [-1, 0], 'color': 'rgba(248,113,113,0.25)'},
                    {'range': [0, 0.5], 'color': 'rgba(250,204,21,0.25)'},
                    {'range': [0.5, 1], 'color': 'rgba(74,222,128,0.25)'},
                ],
            }
        ))
        gauge.update_layout(template="plotly_dark", paper_bgcolor="rgba(0,0,0,0)",
                             margin=dict(l=10, r=10, t=30, b=10), height=380)
        st.plotly_chart(gauge, use_container_width=True)
 
    st.markdown("#### 🎨 Interactive Cluster Plot")
    scatter_fig = px.scatter(
        df, x=primary_feature, y=secondary_feature,
        color="Customer_Type",
        color_discrete_map=BADGE_COLORS,
        size=[10] * len(df),
        hover_data=df.columns,
        opacity=0.85
    )
    scatter_fig.update_layout(
        template="plotly_dark", paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
        legend_title_text="Customer Type", height=480, margin=dict(l=10, r=10, t=30, b=10)
    )
    st.plotly_chart(scatter_fig, use_container_width=True)
 
    st.markdown("#### 📋 Cluster Summary")
    display_summary = summary.copy()
    display_summary["Customer_Type"] = [customer_type_map.get(i, "Noise / Outlier") for i in display_summary.index]
    display_summary["Count"] = valid_df["Cluster"].value_counts().reindex(display_summary.index)
    st.dataframe(display_summary, use_container_width=True)
    if noise_present:
        st.caption(f"⚪ {int((df['Cluster']==-1).sum())} points classified as Noise/Outlier by DBSCAN "
                   f"(not shown in the summary table above).")
 
# ---------------------------------------------------------
# TAB 2 — SEARCH + EXPLORE
# ---------------------------------------------------------
with tab2:
    st.markdown("#### 🔍 Search Customer")
    search_term = st.text_input("Search by ID / Name / row value")
 
    display_df = filtered_df.copy()
    display_df["Segment"] = display_df["Customer_Type"].apply(badge_html)
 
    if search_term:
        if id_col == "(row index)":
            try:
                display_df = display_df.loc[[int(search_term)]]
            except (ValueError, KeyError):
                display_df = display_df.iloc[0:0]
        else:
            display_df = display_df[
                display_df[id_col].astype(str).str.contains(search_term, case=False, na=False)
            ]
 
    st.markdown(f"**{len(display_df)}** customer(s) found")
    st.write(display_df.drop(columns=["Segment"]).to_html(escape=False, index=False) +
             "<br>", unsafe_allow_html=True)
 
    st.markdown("##### Segment Preview")
    for _, row in display_df.head(20).iterrows():
        label = row[id_col] if id_col != "(row index)" else f"Row {row.name}"
 
        def fmt(val, is_num):
            return f"{val:.1f}" if is_num and isinstance(val, (int, float, np.floating)) else f"{val}"
 
        st.markdown(
            f"{badge_html(row['Customer_Type'])} &nbsp; **{label}** — "
            f"{primary_feature}: {fmt(row[primary_feature], primary_is_numeric)}, "
            f"{secondary_feature}: {fmt(row[secondary_feature], secondary_is_numeric)}",
            unsafe_allow_html=True
        )
 
# ---------------------------------------------------------
# TAB 3 — SINGLE CUSTOMER PREDICTION
# ---------------------------------------------------------
with tab3:
    st.markdown("#### 👤 Predict a New Customer's Segment")
    st.write("Enter values for a new or hypothetical customer to see which segment they fall into.")
    if algorithm not in ("KMeans", "Gaussian Mixture"):
        st.caption(f"ℹ️ {algorithm} has no native predict function — this uses a nearest-cluster-centroid "
                   f"approximation instead.")
 
    input_cols = st.columns(len(features))
    input_vals = {}
    for i, f in enumerate(features):
        with input_cols[i]:
            if f in numeric_features:
                input_vals[f] = st.number_input(f, value=float(df[f].mean()), step=1.0, format="%.2f")
            else:
                options = sorted(df[f].dropna().unique().tolist(), key=lambda x: str(x))
                input_vals[f] = st.selectbox(f, options, index=0)
 
    if st.button("🔮 Predict Segment"):
        new_point = pd.DataFrame([input_vals])[features]
        new_scaled = preprocessor.transform(new_point)
        if hasattr(new_scaled, "toarray"):
            new_scaled = new_scaled.toarray()
        pred_cluster = predict_new_point(new_scaled)
        pred_type = customer_type_map.get(pred_cluster, "Noise / Outlier")
 
        st.markdown(f"### Result: {badge_html(pred_type)}", unsafe_allow_html=True)
 
        if pred_type == "High Value Customer":
            st.success("This customer belongs to your top segment — prioritize retention & loyalty offers! 🟢")
        elif pred_type == "Medium Value Customer":
            st.warning("This customer has growth potential — consider upsell campaigns. 🟡")
        elif pred_type == "Low Value Customer":
            st.error("This customer is low value — consider engagement or win-back campaigns. 🔴")
        else:
            st.info("This customer doesn't fit cleanly into an existing segment — worth a manual review. ⚪")
 
# ---------------------------------------------------------
# TAB 4 — DATA TABLE + DOWNLOADS
# ---------------------------------------------------------
with tab4:
    st.markdown("#### 📋 Filtered Cluster Data")
    st.dataframe(filtered_df, use_container_width=True, height=420)
 
    dl1, dl2 = st.columns(2)
    with dl1:
        csv_bytes = filtered_df.to_csv(index=False).encode("utf-8")
        st.download_button("📥 Download CSV", data=csv_bytes,
                            file_name="clustered_customers.csv", mime="text/csv")
    with dl2:
        excel_bytes = to_excel_bytes(filtered_df)
        st.download_button("📥 Download Excel", data=excel_bytes,
                            file_name="clustered_customers.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
 
# ---------------------------------------------------------
# TAB 5 — BUSINESS INSIGHTS (deeper)
# ---------------------------------------------------------
with tab5:
    st.markdown("#### 📊 Auto-Generated Business Insights")
    insight_lines = generate_insight_lines()
    for line in insight_lines:
        st.markdown(f"<div class='insight-box'>💡 {line}</div>", unsafe_allow_html=True)
 
    st.markdown("##### Segment Distribution")
    pie_fig = px.pie(df, names="Customer_Type", hole=0.55,
                      color="Customer_Type", color_discrete_map=BADGE_COLORS)
    pie_fig.update_layout(template="plotly_dark", paper_bgcolor="rgba(0,0,0,0)",
                           margin=dict(l=10, r=10, t=10, b=10), height=380)
    st.plotly_chart(pie_fig, use_container_width=True)
 
    if numeric_features:
        st.markdown("##### Average Feature Value by Segment")
        bar_fig = px.bar(
            summary.reset_index().assign(Customer_Type=lambda d: d["Cluster"].map(customer_type_map)),
            x="Customer_Type", y=numeric_features, barmode="group",
            color_discrete_sequence=["#a78bfa", "#38bdf8", "#34d399", "#facc15"]
        )
        bar_fig.update_layout(template="plotly_dark", paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
                               margin=dict(l=10, r=10, t=10, b=10), height=380)
        st.plotly_chart(bar_fig, use_container_width=True)
 
    if len(numeric_features) >= 2:
        st.markdown("##### 🔗 Correlation Between Numeric Features")
        corr = df[numeric_features].corr().round(2)
        corr_fig = px.imshow(corr, text_auto=True, color_continuous_scale="Purples",
                              aspect="auto")
        corr_fig.update_layout(template="plotly_dark", paper_bgcolor="rgba(0,0,0,0)",
                                margin=dict(l=10, r=10, t=10, b=10), height=380)
        st.plotly_chart(corr_fig, use_container_width=True)
 
    st.markdown("---")
    st.markdown("##### ✨ AI Executive Summary (optional)")
    st.caption("Uses your own Anthropic API key — it's only used for this session and never stored.")
    if not ANTHROPIC_OK:
        st.info("Install the `anthropic` package (`pip install anthropic`) to enable this feature.")
    else:
        api_key = st.text_input("Anthropic API Key", type="password", key="anthropic_key_tab5")
        if st.button("✨ Generate AI Executive Summary"):
            if not api_key:
                st.warning("Please enter an API key first.")
            else:
                try:
                    client = anthropic.Anthropic(api_key=api_key)
                    context = (
                        f"Dataset: {len(df)} customers. Features used: {', '.join(features)}. "
                        f"Algorithm: {algorithm}. Silhouette score: {score:.3f}.\n"
                        f"Cluster summary table:\n{summary.to_string()}\n\n"
                        f"Rule-based findings:\n- " + "\n- ".join(insight_lines)
                    )
                    resp = client.messages.create(
                        model="claude-sonnet-5",
                        max_tokens=500,
                        messages=[{
                            "role": "user",
                            "content": f"You are a business analyst. Based on this customer segmentation "
                                       f"data, write a concise (150-200 word) executive summary with 2-3 "
                                       f"concrete action recommendations:\n\n{context}"
                        }]
                    )
                    ai_text = "".join([b.text for b in resp.content if hasattr(b, "text")])
                    st.session_state["ai_summary"] = ai_text
                except Exception as e:
                    st.error(f"Could not generate summary: {e}")
        if st.session_state.get("ai_summary"):
            st.markdown(f"<div class='insight-box'>🤖 {st.session_state['ai_summary']}</div>",
                        unsafe_allow_html=True)
 
# ---------------------------------------------------------
# TAB 6 — REPORTS (PDF / PPTX export)
# ---------------------------------------------------------
with tab6:
    st.markdown("#### 📄 Export a Shareable Report")
    st.write("Generates a report with your KPIs, cluster summary, and insights (AI summary included if generated above).")
 
    def make_pie_png():
        fig, ax = plt.subplots(figsize=(4, 4))
        counts = df["Customer_Type"].value_counts()
        colors = [BADGE_COLORS.get(c, "#a78bfa") for c in counts.index]
        ax.pie(counts.values, labels=counts.index, autopct="%1.1f%%", colors=colors,
               textprops={"fontsize": 8})
        buf = BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
        plt.close(fig)
        buf.seek(0)
        return buf
 
    def make_bar_png():
        if not numeric_features:
            return None
        fig, ax = plt.subplots(figsize=(6, 4))
        plot_df = summary.reset_index().assign(Customer_Type=lambda d: d["Cluster"].map(customer_type_map))
        plot_df.plot(x="Customer_Type", y=numeric_features, kind="bar", ax=ax)
        ax.set_ylabel("Average value")
        buf = BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
        plt.close(fig)
        buf.seek(0)
        return buf
 
    insight_lines = generate_insight_lines()
    ai_summary_text = st.session_state.get("ai_summary", "")
 
    col_pdf, col_ppt = st.columns(2)
 
    with col_pdf:
        st.markdown("**PDF Report**")
        if not (REPORTLAB_OK and MATPLOTLIB_OK):
            st.info("Install `reportlab` and `matplotlib` (`pip install reportlab matplotlib`) to enable PDF export.")
        else:
            if st.button("📥 Build PDF Report"):
                buf = BytesIO()
                doc = SimpleDocTemplate(buf, pagesize=A4)
                styles = getSampleStyleSheet()
                elements = [
                    Paragraph("Customer Segmentation Report", styles["Title"]),
                    Paragraph(f"Generated on {pd.Timestamp.now().strftime('%d %b %Y, %H:%M')}", styles["Normal"]),
                    Spacer(1, 14),
                ]
                kpi_data = [["Metric", "Value"],
                            ["Total Customers", f"{len(df)}"],
                            ["Algorithm", algorithm],
                            ["Silhouette Score", f"{score:.3f}"],
                            ["High Value %", f"{(df['Customer_Type']=='High Value Customer').mean()*100:.1f}%"]]
                t = Table(kpi_data, hAlign="LEFT")
                t.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), rl_colors.HexColor("#7c3aed")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), rl_colors.white),
                    ("GRID", (0, 0), (-1, -1), 0.5, rl_colors.grey),
                    ("FONTSIZE", (0, 0), (-1, -1), 9),
                ]))
                elements += [t, Spacer(1, 16), Paragraph("Key Insights", styles["Heading2"])]
                for line in insight_lines:
                    elements.append(Paragraph("• " + line, styles["Normal"]))
                if ai_summary_text:
                    elements += [Spacer(1, 12), Paragraph("AI Executive Summary", styles["Heading2"]),
                                 Paragraph(ai_summary_text, styles["Normal"])]
                elements += [Spacer(1, 16), Paragraph("Segment Distribution", styles["Heading2"]),
                             RLImage(make_pie_png(), width=9 * cm, height=9 * cm)]
                bar_png = make_bar_png()
                if bar_png:
                    elements += [Spacer(1, 12), Paragraph("Average Feature Value by Segment", styles["Heading2"]),
                                 RLImage(bar_png, width=14 * cm, height=9 * cm)]
                doc.build(elements)
                buf.seek(0)
                st.download_button("⬇️ Download PDF", data=buf, file_name="segmentation_report.pdf",
                                    mime="application/pdf")
 
    with col_ppt:
        st.markdown("**PowerPoint Report**")
        if not (PPTX_OK and MATPLOTLIB_OK):
            st.info("Install `python-pptx` and `matplotlib` (`pip install python-pptx matplotlib`) to enable PPTX export.")
        else:
            if st.button("📥 Build PPTX Report"):
                prs = Presentation()
                blank = prs.slide_layouts[6]
 
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = "Customer Segmentation Report"
                slide.placeholders[1].text = pd.Timestamp.now().strftime("%d %b %Y")
 
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Key Metrics"
                body = slide.placeholders[1].text_frame
                body.text = f"Total Customers: {len(df)}"
                for txt in [f"Algorithm: {algorithm}", f"Silhouette Score: {score:.3f}",
                            f"High Value %: {(df['Customer_Type']=='High Value Customer').mean()*100:.1f}%"]:
                    p = body.add_paragraph()
                    p.text = txt
 
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Key Insights"
                body = slide.placeholders[1].text_frame
                body.text = insight_lines[0] if insight_lines else ""
                for line in insight_lines[1:6]:
                    p = body.add_paragraph()
                    p.text = line
 
                if ai_summary_text:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = "AI Executive Summary"
                    slide.placeholders[1].text_frame.text = ai_summary_text
 
                slide = prs.slides.add_slide(blank)
                slide.shapes.add_picture(make_pie_png(), Inches(2), Inches(1), height=Inches(5))
 
                bar_png = make_bar_png()
                if bar_png:
                    slide = prs.slides.add_slide(blank)
                    slide.shapes.add_picture(bar_png, Inches(1), Inches(1), width=Inches(8))
 
                buf = BytesIO()
                prs.save(buf)
                buf.seek(0)
                st.download_button("⬇️ Download PPTX", data=buf, file_name="segmentation_report.pptx",
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
 
st.markdown("<br><center style='color:#5f6180; font-size:0.8rem;'>Built with ❤️ using Streamlit • Multi-Algorithm Customer Segmentation</center>", unsafe_allow_html=True)
